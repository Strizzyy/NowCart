from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand colors ──────────────────────────────────────────────
GREEN       = RGBColor(0x3b, 0xb7, 0x7e)
GREEN_DARK  = RGBColor(0x15, 0x73, 0x47)
GREEN_LIGHT = RGBColor(0xde, 0xf9, 0xec)
AMBER       = RGBColor(0xf5, 0xa6, 0x23)
AMBER_DARK  = RGBColor(0xb9, 0x74, 0x0a)
ACCENT      = RGBColor(0xe2, 0x3d, 0x6d)
DARK        = RGBColor(0x1d, 0x32, 0x42)
MUTED       = RGBColor(0x5f, 0x6b, 0x66)
FAINT       = RGBColor(0x8a, 0x93, 0x8e)
SURFACE     = RGBColor(0xff, 0xff, 0xff)
LIGHT_BG    = RGBColor(0xf4, 0xf6, 0xfa)
BLUE        = RGBColor(0x1d, 0x4e, 0xd8)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]   # completely blank layout

# ── Helper functions ──────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(blank)

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, x, y, w, h, bg_color=None, border_color=None, border_w=Pt(1)):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    if bg_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_w
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, x, y, w, h, size=14, bold=False, color=None,
        align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color or DARK
    return tb

def heading(slide, text, x=0.35, y=0.15, w=12.6, size=28, color=None):
    return txt(slide, text, x, y, w, 0.55, size=size, bold=True,
               color=color or DARK, align=PP_ALIGN.LEFT)

def slide_label(slide, label):
    txt(slide, label, 0.35, 0.05, 12, 0.18, size=8, color=FAINT)

def divider(slide, y, color=None):
    ln = slide.shapes.add_connector(1,
         Inches(0.35), Inches(y), Inches(12.98), Inches(y))
    ln.line.color.rgb = color or RGBColor(0xe3, 0xe7, 0xe5)
    ln.line.width = Pt(1)

def green_bar(slide):
    """Top accent bar"""
    b = box(slide, 0, 0, 13.33, 0.07, bg_color=GREEN)

def card_box(slide, x, y, w, h, accent=None):
    """White card with optional left accent bar"""
    b = box(slide, x, y, w, h, bg_color=SURFACE,
            border_color=RGBColor(0xe3,0xe7,0xe5), border_w=Pt(1))
    if accent:
        box(slide, x, y, 0.05, h, bg_color=accent)
    return b

# ══════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════
s1 = add_slide()
bg(s1, LIGHT_BG)
green_bar(s1)

# Tags row
txt(s1,"48-HOUR HACKATHON  ·  AMAZON NOW TRACK  ·  TEAM CODYSSEY",
    0.35,0.18,12,0.25, size=9, bold=True, color=FAINT)

# NowCart wordmark
txt(s1,"Now", 0.35,0.55,2.2,0.85, size=60, bold=True, color=GREEN_DARK)
txt(s1,"Cart  🛒", 2.45,0.55,3.8,0.85, size=60, bold=True, color=DARK)

# Tagline
txt(s1,"Reimagining Urgent Shopping", 0.35,1.45,8,0.38,
    size=20, bold=False, color=MUTED, italic=True)
txt(s1,'"Quick commerce solved delivery.  We solve the deciding."',
    0.35,1.85,9,0.38, size=16, bold=True, color=GREEN_DARK, italic=True)

divider(s1, 2.35)

# Event details
txt(s1,"HackOn with Amazon  ·  Amazon Now Track  ·  15 / 06 / 2026",
    0.35,2.42,8,0.28, size=10, color=FAINT)

# Team card (right side)
card_box(s1, 9.1, 0.5, 3.9, 3.2, accent=GREEN)
txt(s1,"TEAM CODYSSEY", 9.22,0.58,3.6,0.22, size=8, bold=True, color=FAINT)
members = [
    ("👨‍💻  Rohan Singh","RGIPT, Bangalore  ·  Backend + Deployment"),
    ("👨‍💻  Anuj Kumar Yadav","NSUT, Delhi  ·  Frontend + Backend"),
    ("👨‍🔬  Baibhav Kundu","NSUT, Delhi  ·  Machine Learning"),
]
for i,(name,role) in enumerate(members):
    yy = 0.90 + i*0.72
    txt(s1, name, 9.22, yy,      3.6,0.28, size=12, bold=True,  color=DARK)
    txt(s1, role, 9.22, yy+0.28, 3.6,0.28, size=10, bold=False, color=MUTED)

# Output box
card_box(s1, 9.1,3.82,3.9,0.98)
txt(s1,"Need  →  Ready Cart", 9.22,3.90,3.7,0.35,
    size=14, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)
txt(s1,"5 front doors  ·  1 reasoning engine  ·  0 dead ends",
    9.22,4.26,3.7,0.28, size=9, color=MUTED, align=PP_ALIGN.CENTER)

# Speaker note placeholder
txt(s1,'Speaker note: "We\'re Team Codyssey — quick commerce solved delivery, but nobody solved the 5 minutes before that. We do."',
    0.35,6.95,12.6,0.45, size=8, color=FAINT, italic=True)

# ══════════════════════════════════════════════════════════════
# SLIDE 2 — SOLUTION OVERVIEW
# ══════════════════════════════════════════════════════════════
s2 = add_slide()
bg(s2, LIGHT_BG)
green_bar(s2)
slide_label(s2,"Slide 2  ·  Solution Overview  ·  0:20 – 1:00")
heading(s2,"NowCart: The Requirement-First Commerce Layer")
txt(s2,"The world's first shopping engine that starts with your need — not a search box.",
    0.35,0.72,12.6,0.28, size=12, color=MUTED)
divider(s2,1.08)

# Left — comparison cards
card_box(s2,0.35,1.18,5.8,0.85, accent=BLUE)
txt(s2,"EVERY PLATFORM TODAY",0.48,1.22,5.5,0.22, size=8, bold=True, color=FAINT)
txt(s2,"Product → Search → Browse → Add → Repeat",
    0.48,1.42,5.5,0.28, size=12, bold=False, color=MUTED, italic=True)
txt(s2,"5–7 minutes. Every. Single. Order.",
    0.48,1.70,5.5,0.22, size=9, color=FAINT)

card_box(s2,0.35,2.18,5.8,1.05, accent=GREEN)
txt(s2,"NOWCART",0.48,2.22,5.5,0.22, size=8, bold=True, color=GREEN_DARK)
txt(s2,"Requirement  →  Ready Cart",
    0.48,2.42,5.5,0.38, size=18, bold=True, color=DARK)
txt(s2,'"Penne pasta for 2" → full ingredient cart in seconds.',
    0.48,2.80,5.5,0.28, size=11, color=MUTED)

# Right — 3 value cards
props = [
    (GREEN,"🎯 What is NowCart?",
     "A universal reasoning layer on top of any quick commerce catalog. Turns a high-level need into a checkout-ready cart automatically."),
    (AMBER,"🔄 The Core Reversal",
     "We don't solve delivery times. We solve decision fatigue. Intent is the input. Cart is the output."),
    (ACCENT,"✅ Informed Choice, Not Forced Choice",
     "Transparent ranking. NowCart Verified badge = most ordered + highest rated. Economical Alt with exact saving shown. No brand bias."),
]
for i,(accent,title,body) in enumerate(props):
    yy = 1.18 + i*1.05
    card_box(s2,6.45,yy,6.5,0.95, accent=accent)
    txt(s2,title,6.58,yy+0.08,6.2,0.28, size=12, bold=True, color=DARK)
    txt(s2,body, 6.58,yy+0.36,6.2,0.52, size=10, color=MUTED)

txt(s2,'Speaker note: "NowCart is a universal reasoning layer — express a need, get a checkout-ready cart. We reverse the entire flow from product-first to requirement-first."',
    0.35,6.95,12.6,0.45, size=8, color=FAINT, italic=True)

# ══════════════════════════════════════════════════════════════
# SLIDE 3 — PROBLEM + CUSTOMER PAIN POINT
# ══════════════════════════════════════════════════════════════
s3 = add_slide()
bg(s3, LIGHT_BG)
green_bar(s3)
slide_label(s3,"Slide 3  ·  The Bottleneck  ·  1:00 – 1:50")
heading(s3,"The Last Unsolved Problem in Quick Commerce")
divider(s3,0.82)

# Pain point scenario strip
card_box(s3,0.35,0.90,12.6,1.12, accent=ACCENT)
txt(s3,"😤  A Real User's Thursday Evening",
    0.50,0.94,12,0.28, size=11, bold=True, color=ACCENT)
txt(s3,"Priya wants to cook Paneer Butter Masala. She opens Blinkit, searches \"paneer\" — 12 results, 4 pack sizes, 3 brands. Picks one. Searches \"tomato puree\" — out of stock. Searches again. Then butter, cream, spices — one by one. 7 minutes later she's still building the cart. Delivery: 10 min. Deciding: just as long.",
    0.50,1.22,12,0.72, size=10, color=DARK)

# 3 stat boxes
stats = [
    (GREEN_DARK,   "5–7 min","wasted per order across 300M+ Indian shoppers\n(RedSeer, 2025)"),
    (ACCENT,       "~70%",   "of online carts abandoned before checkout — 7 of 10\n(Dynamic Yield, 2025)"),
    (BLUE,         "$7B+",   "quick-commerce market in India — delivery perfected,\ncart-building unsolved  (Mordor Intel, 2025)"),
]
for i,(color,num,lbl) in enumerate(stats):
    xx = 0.35 + i*4.33
    card_box(s3,xx,2.18,4.1,0.88)
    txt(s3,num, xx+0.12,2.24,3.8,0.42, size=26, bold=True, color=color, align=PP_ALIGN.CENTER)
    txt(s3,lbl, xx+0.12,2.62,3.8,0.38, size=9, color=MUTED, align=PP_ALIGN.CENTER)

# 3 numbered problem points (left)
problems = [
    ("Cart-building is the hidden bottleneck",
     "Every app forces item-by-item search, 4–6 variant comparisons per product, and manual OOS hunting. Delivery is 10 min. Building the cart takes the same."),
    ("OOS dead-ends drive permanent churn",
     "User hits an OOS item mid-cart — app shows nothing. They switch to a rival app. Every dead-end is a permanent trust loss, not a temporary inconvenience."),
    ("Paradox of choice — no decision layer",
     "6 variants of tomato puree: different brands, weights, prices. No platform tells you which fits your recipe. Users freeze and abandon."),
]
for i,(title,body) in enumerate(problems):
    yy = 3.22 + i*1.02
    card_box(s3,0.35,yy,7.0,0.94, accent=GREEN)
    txt(s3,str(i+1), 0.48,yy+0.08,0.35,0.35, size=14, bold=True, color=SURFACE, align=PP_ALIGN.CENTER)
    box(s3,0.47,yy+0.07,0.36,0.36, bg_color=GREEN, border_color=GREEN)
    txt(s3,title, 0.95,yy+0.06,6.2,0.26, size=12, bold=True, color=DARK)
    txt(s3,body,  0.95,yy+0.34,6.2,0.52, size=10, color=MUTED)

# Gap card (right)
card_box(s3,7.55,3.22,5.4,3.06, accent=GREEN)
txt(s3,"THE GAP NOBODY FILLS", 7.70,3.28,5.1,0.22, size=9, bold=True, color=GREEN_DARK)
divider(s3,3.56)
rows = [
    ("🚀","10-min delivery  ✅  Solved","Blinkit, Zepto, Amazon Fresh, Instamart",GREEN_DARK),
    ("🧠","10-min cart-building  ❌  Unsolved","No platform has built this. Until NowCart.",ACCENT),
    ("💬",'"I know what I want to cook. I just don\'t want to\nspend 7 minutes searching for every ingredient."',"— every household shopper, every time",MUTED),
]
for i,(icon,line1,line2,c) in enumerate(rows):
    yy = 3.65 + i*0.88
    txt(s3,icon,  7.65,yy,0.45,0.38, size=18)
    txt(s3,line1, 8.12,yy,4.65,0.30, size=11, bold=True, color=c)
    txt(s3,line2, 8.12,yy+0.30,4.65,0.38, size=9, color=MUTED, italic=(i==2))

txt(s3,'Speaker note: "Imagine Priya — she knows what to cook but spends 7 minutes building the cart. Delivery: 10 min. Deciding: just as long. 300 million users, every day. NowCart closes that gap."',
    0.35,6.95,12.6,0.45, size=8, color=FAINT, italic=True)

# ══════════════════════════════════════════════════════════════
# SLIDE 4 — CUSTOMER FIT + NOVELTY
# ══════════════════════════════════════════════════════════════
s4 = add_slide()
bg(s4, LIGHT_BG)
green_bar(s4)
slide_label(s4,"Slide 4  ·  Customer Fit & Novelty  ·  1:50 – 3:00")
heading(s4,"Built for Real People. Designed Around Intent.")
divider(s4,0.82)

# Table header
cols = [1.5,2.5,2.8,2.8]
headers = ["User Group","Real Pain Point","What Happens Today","NowCart's Answer"]
xx = 0.35
for i,(w,h) in enumerate(zip(cols,headers)):
    box(s4,xx,0.90,w,0.32, bg_color=DARK, border_color=DARK)
    txt(s4,h, xx+0.06,0.92,w-0.1,0.26, size=9, bold=True, color=SURFACE)
    xx += w

# Table rows
rows = [
    ("🧑‍🍳 Kitchen Novices\n(saw a dish online)",
     "Wants to cook a dish from Instagram. No idea what to buy, in what qty, or which brand.",
     "Googles recipe, manually searches each ingredient, gives up halfway.",
     "Show + Share: Snap dish or paste YouTube link → full cart in seconds"),
    ("👴 Elderly / Voice-First\n(non-tech users)",
     "Wants groceries but can't navigate text-heavy apps — small buttons, typing are real barriers.",
     "Asks family member to order. Dependent, slow, not independent.",
     "Speak: Say 'dal chawal for 2' → cart appears. Zero typing, zero dependency."),
    ("🏠 Daily Household Shopper\n(weekly planner)",
     "Knows what to cook but wastes 5–7 min, hits OOS dead-ends, compares 6 variants per item.",
     "Cart abandoned 70% of the time — session ends in frustration, not checkout.",
     "Constrain + Subscribe: Budget cart or auto-predicted restock, ready before they ask."),
]
accents = [GREEN, AMBER, BLUE]
for ri,(r0,r1,r2,r3) in enumerate(rows):
    yy = 1.25 + ri*1.25
    bg_c = SURFACE if ri%2==0 else LIGHT_BG
    xx = 0.35
    for ci,(w,cell) in enumerate(zip(cols,[r0,r1,r2,r3])):
        box(s4,xx,yy,w,1.18, bg_color=bg_c,
            border_color=RGBColor(0xe3,0xe7,0xe5), border_w=Pt(1))
        color = accents[ri] if ci==0 else (ACCENT if ci==2 else DARK)
        txt(s4,cell, xx+0.08,yy+0.06,w-0.14,1.05, size=9,
            bold=(ci==0), color=color)
        xx += w

# Bottom highlight strip
card_box(s4,0.35,4.95,12.6,0.52, accent=GREEN)
txt(s4,"What they all need: Not a better search box — an interface that takes a raw need and hands back a ready cart. One missing item never breaks the flow. It's swapped with the next best match, transparently.",
    0.50,5.01,12.2,0.40, size=10, color=DARK)

# 3 novelty cards
novelties = [
    (GREEN,"① Unique Angle — Intent → Cart",
     "Only platform that goes intent→cart, not product→you decide. 'Penne pasta for 2' builds the cart."),
    (AMBER,"② Informed Choice, Not Forced Choice",
     "NowCart Verified badge = most ordered + highest rated. Economical Alt always shown. No hidden AI bias."),
    (BLUE,"③ No Single Point of Failure",
     "Photo, link, voice, budget, subscription all feed one reasoning framework. Step fails? Pipeline skips it."),
]
for i,(accent,title,body) in enumerate(novelties):
    xx = 0.35 + i*4.33
    card_box(s4,xx,5.58,4.1,0.82, accent=accent)
    txt(s4,title, xx+0.14,5.62,3.8,0.26, size=10, bold=True, color=DARK)
    txt(s4,body,  xx+0.14,5.90,3.8,0.44, size=9, color=MUTED)

txt(s4,'Speaker note: "Three users, one problem. Kitchen novice gives up. Elderly user calls their kid. Household shopper abandons the cart. None need a better search box — they need intent-in, cart-out."',
    0.35,6.95,12.6,0.45, size=8, color=FAINT, italic=True)

# ══════════════════════════════════════════════════════════════
# SLIDE 5 — FIVE FRONT DOORS
# ══════════════════════════════════════════════════════════════
s5 = add_slide()
bg(s5, LIGHT_BG)
green_bar(s5)
slide_label(s5,"Slide 5  ·  Five Front Doors  ·  3:00 – 3:50")
heading(s5,"Five Ways In.  One Confident Cart Out.")
divider(s5,0.82)

doors = [
    (GREEN,  "📸","SHOW",   "Snap a dish photo → Gemini Vision extracts ingredients → LangGraph builds the cart"),
    (BLUE,   "🔗","SHARE",  "Paste any recipe or YouTube URL → LLM parses ingredients → cart built"),
    (AMBER,  "🎤","SPEAK",  "Say 'biryani for 4' → semantic intent match + weight resolution → confidence-scored cart"),
    (ACCENT, "💰","CONSTRAIN","Set ₹500 budget → full pipeline runs → greedy knapsack trims by confidence until total fits"),
    (RGBColor(0x7c,0x3a,0xed),"🔁","SUBSCRIBE","Order history → inter-purchase interval analysis → restock cart pre-built before you ask"),
]
dw = 2.36
for i,(accent,icon,label,desc) in enumerate(doors):
    xx = 0.35 + i*(dw+0.12)
    card_box(s5,xx,0.95,dw,2.45)
    box(s5,xx,0.95,dw,0.06, bg_color=accent)
    txt(s5,icon,  xx+0.1,1.08,dw-0.2,0.48, size=28, align=PP_ALIGN.CENTER)
    txt(s5,label, xx+0.1,1.58,dw-0.2,0.28, size=12, bold=True, color=accent, align=PP_ALIGN.CENTER)
    txt(s5,desc,  xx+0.1,1.92,dw-0.2,1.35, size=9, color=MUTED, align=PP_ALIGN.CENTER)

# Convergence arrow
txt(s5,"▼  all five feed the same engine  ▼",
    0.35,3.58,12.6,0.30, size=10, color=FAINT, align=PP_ALIGN.CENTER)

# Engine box
card_box(s5,2.65,3.95,8.0,1.02, accent=GREEN)
txt(s5,"Outcome Engine  (LangGraph Brain)",
    2.78,3.99,7.7,0.28, size=12, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)
txt(s5,"decompose  →  match  →  semantic re-rank  →  score  →  resolve OOS",
    2.78,4.28,7.7,0.28, size=11, bold=False, color=DARK, align=PP_ALIGN.CENTER)

txt(s5,"▼",4.35,5.06,4.6,0.24, size=16, color=FAINT, align=PP_ALIGN.CENTER)

# Output box
card_box(s5,2.65,5.30,8.0,0.90, accent=GREEN_DARK)
txt(s5,"One Confident Cart",
    2.78,5.34,7.7,0.28, size=13, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)
txt(s5,"✓ Recommended Pick    |    💰 Economical Alt    |    📊 Confidence %    |    💬 One-line reason",
    2.78,5.62,7.7,0.28, size=10, color=DARK, align=PP_ALIGN.CENTER)

txt(s5,'Speaker note: "Five ways to express a need — photo, voice, link, budget, or predictive restock. All five feed the same engine. Output is always one confident cart."',
    0.35,6.95,12.6,0.45, size=8, color=FAINT, italic=True)

# ══════════════════════════════════════════════════════════════
# SLIDE 6 — WORKFLOW  Express · Engine · Confirm
# ══════════════════════════════════════════════════════════════
s6 = add_slide()
bg(s6, LIGHT_BG)
green_bar(s6)
slide_label(s6,"Slide 6  ·  User Workflow  ·  3:50 – 4:20  ·  Then → Demo")
heading(s6,"Express  ·  Engine  ·  Confirm")
txt(s6,"Three steps from need to checkout. No search. No dead ends. No wasted minutes.",
    0.35,0.72,12.6,0.28, size=12, color=MUTED)
divider(s6,1.08)

steps = [
    (GREEN,  "✏️","1. EXPRESS",
     "📸  Snap a dish\n🔗  Paste a recipe link\n🎤  Speak a meal\n💰  Set a budget\n🔁  Let it predict"),
    (AMBER,  "⚙️","2. ENGINE",
     "Decomposes your intent\nMatches 9,534 products\nBi + Cross + Fuzzy retrieval\nResolves out-of-stock\nScores confidence per item"),
    (GREEN_DARK,"✅","3. CONFIRM",
     "One cart with confidence scores\nTransparent substitutions shown\nRefine conversationally:\n'make it vegan' / 'remove onion'\nCheckout in 2 taps"),
]
sw = 3.8
for i,(accent,icon,title,body) in enumerate(steps):
    xx = 0.35 + i*(sw+0.46)
    card_box(s6,xx,1.18,sw,4.55, accent=accent)
    txt(s6,icon,  xx+0.1,1.30,sw-0.2,0.60, size=36, align=PP_ALIGN.CENTER)
    txt(s6,title, xx+0.1,1.96,sw-0.2,0.38, size=16, bold=True, color=accent, align=PP_ALIGN.CENTER)
    txt(s6,body,  xx+0.1,2.40,sw-0.2,3.20, size=11, color=MUTED, align=PP_ALIGN.CENTER)
    if i < 2:
        txt(s6,"▶", xx+sw+0.05,3.12,0.38,0.38, size=20, color=FAINT, align=PP_ALIGN.CENTER)

# Refinement strip
card_box(s6,0.35,5.85,12.6,0.50, accent=GREEN)
txt(s6,'Conversational refinement after cart is built:   "make it cheaper"   ·   "I\'m vegetarian"   ·   "remove onion"   ·   "reduce to ₹300"',
    0.50,5.91,12.2,0.38, size=11, bold=False, color=DARK, align=PP_ALIGN.CENTER)

txt(s6,'Speaker note: "Express your need. Engine decomposes, matches, resolves OOS, scores picks. One confident cart. Refine conversationally. Checkout in 2 taps. Let\'s see this live."',
    0.35,6.95,12.6,0.45, size=8, color=FAINT, italic=True)

# ══════════════════════════════════════════════════════════════
# SLIDE 7 — DEMO
# ══════════════════════════════════════════════════════════════
s7 = add_slide()
bg(s7, LIGHT_BG)
green_bar(s7)
slide_label(s7,"Slide 7  ·  Live Prototype  ·  4:20 – 7:50")
heading(s7,"Live Prototype Demonstration", size=28)
divider(s7,0.82)

# Browser chrome mockup
card_box(s7,1.5,0.95,10.3,5.35)
box(s7,1.5,0.95,10.3,0.38, bg_color=LIGHT_BG, border_color=RGBColor(0xe3,0xe7,0xe5))
# traffic lights
for ci,c in enumerate([ACCENT, AMBER, GREEN]):
    box(s7,1.68+ci*0.28,1.05,0.18,0.18, bg_color=c)
txt(s7,"nowcart.live",3.0,1.02,2.5,0.22, size=9, color=FAINT)

txt(s7,"🎬",5.6,1.85,2.1,0.90, size=56, align=PP_ALIGN.CENTER)
txt(s7,"EMBED DEMO VIDEO HERE",
    2.5,2.80,8.3,0.40, size=18, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
txt(s7,"3:30 – 3:40 min walkthrough",
    2.5,3.22,8.3,0.28, size=12, color=FAINT, align=PP_ALIGN.CENTER)

features = ["📸 Show: Dish Photo","🔗 Share: YouTube Link","🎤 Speak: Biryani for 4","💰 Constrain: ₹500 Budget","🔁 Subscribe: Restock"]
fw = 2.22
for i,f in enumerate(features):
    xx = 1.72 + i*(fw+0.18)
    card_box(s7,xx,3.72,fw,0.45)
    txt(s7,f, xx+0.08,3.76,fw-0.14,0.34, size=9, color=DARK, align=PP_ALIGN.CENTER)

# Deploy info
card_box(s7,2.5,4.40,8.3,0.72)
txt(s7,"📊  Live on AWS  ·  Test login:  rahul@gmail.com  (any password)  ·  Full pantry + prediction experience",
    2.65,4.52,8.0,0.42, size=11, color=DARK, align=PP_ALIGN.CENTER)

txt(s7,'Speaker note: Narrate during video — "Snap Paneer Butter Masala photo → ingredients extracted. Paste YouTube link → parsed real-time. Voice ₹500 budget → knapsack selects highest-confidence essentials. OOS items swapped, never dead-ends."',
    0.35,6.95,12.6,0.45, size=8, color=FAINT, italic=True)

# ══════════════════════════════════════════════════════════════
# SLIDE 8 — TECH ARCHITECTURE + STACK + ALGORITHMS
# ══════════════════════════════════════════════════════════════
s8 = add_slide()
bg(s8, LIGHT_BG)
green_bar(s8)
slide_label(s8,"Slide 8  ·  Architecture + Stack + Algorithms  ·  8:00 – 9:00")
heading(s8,"How It Actually Works: End-to-End")
divider(s8,0.82)

# Pipeline flow
pipeline = [
    (GREEN,  "① Input",      "5 front doors\nPhoto/Voice\nLink/Budget\nPredict"),
    (AMBER,  "② Security",   "SSRF guard\nRate limiting\nInput sanitize\nSpeech→Text"),
    (BLUE,   "③ Decision",   "Assembly needed?\nYES → engine\nNO → exact\nproduct match"),
    (RGBColor(0x7c,0x3a,0xed),"④ Engine",  "LangGraph 10-node\nBi-encoder search\nCross-encoder rank\nRapidFuzz typos\nOOS auto-swap"),
    (GREEN_DARK,"⑤ Output",  "Verified pick\nEcon alt\nScore +\nReason"),
]
pw = [1.9,1.9,1.9,2.5,1.9]
xx = 0.35
for i,((accent,label,body),w) in enumerate(zip(pipeline,pw)):
    card_box(s8,xx,0.95,w,1.65)
    box(s8,xx,0.95,w,0.06, bg_color=accent)
    txt(s8,label, xx+0.08,1.05,w-0.14,0.26, size=9, bold=True, color=accent)
    txt(s8,body,  xx+0.08,1.34,w-0.14,1.20, size=8.5, color=MUTED)
    if i<4:
        txt(s8,"→", xx+w+0.01,1.58,0.22,0.28, size=14, color=FAINT, align=PP_ALIGN.CENTER)
    xx += w + 0.22

# Stack table
headers_t = ["Layer","Technology","Key Algorithm","What It Enables"]
col_w = [1.4,3.4,2.8,4.7]
ty = 2.78
xx = 0.35
for h,w in zip(headers_t,col_w):
    box(s8,xx,ty,w,0.32, bg_color=DARK, border_color=DARK)
    txt(s8,h, xx+0.06,ty+0.04,w-0.1,0.24, size=8, bold=True, color=SURFACE)
    xx+=w

rows_t = [
    ("Frontend","React 19 + Vite + TailwindCSS","—","Instant panels, works on mid-range phones with patchy network"),
    ("Backend","FastAPI + Pydantic 2 (fully async)","—","Thousands of simultaneous carts, validated end-to-end"),
    ("AI / ML","LangGraph + Groq / Gemini / Claude + TF-IDF","Bi→Cross→RapidFuzz retrieval","Semantic search + re-rank + catches 'malai', 'tomatoe', 'paneer'"),
    ("Database","DynamoDB (pay-per-request)","Coeff. of variation restock scoring","Auto-scales; order history powers Subscribe. No cold start."),
    ("Infra","EC2 + Nginx + S3 + CloudFront + Redis + Lambda + SQS","LLM hash caching 1-hr TTL","100 queries = 1 model call + 99 cache hits. Saves ~800ms/repeat."),
]
for ri,row in enumerate(rows_t):
    yy = 3.12 + ri*0.62
    bg_c = SURFACE if ri%2==0 else LIGHT_BG
    xx = 0.35
    for ci,(cell,w) in enumerate(zip(row,col_w)):
        box(s8,xx,yy,w,0.60, bg_color=bg_c, border_color=RGBColor(0xe3,0xe7,0xe5), border_w=Pt(1))
        txt(s8,cell, xx+0.07,yy+0.06,w-0.12,0.48, size=8, color=DARK)
        xx+=w

# Scaling callout
card_box(s8,7.8,6.26,4.85,0.55, accent=GREEN)
txt(s8,"Scales: 100× Auto Scaling  ·  1000× Multi-region + DAX  ·  AI → Lambda+SQS (1K parallel)",
    7.95,6.30,4.6,0.42, size=8.5, bold=True, color=GREEN_DARK)

txt(s8,'Speaker note: "Input→security→decision→LangGraph engine. Matches 9,534 products with 3 retrieval methods. DynamoDB auto-scales. CloudFront is already global. LLM caching cuts 99% of redundant model calls."',
    0.35,6.95,12.6,0.45, size=8, color=FAINT, italic=True)

# ══════════════════════════════════════════════════════════════
# SLIDE 9 — THEME ALIGNMENT + FUTURE VISION
# ══════════════════════════════════════════════════════════════
s9 = add_slide()
bg(s9, LIGHT_BG)
green_bar(s9)
slide_label(s9,"Slide 9  ·  Theme Alignment + Future Vision  ·  9:00 – 10:00")
heading(s9,"All 3 Amazon Now Areas.  One Engine to Rule Them All.")
divider(s9,0.82)

# Theme 3 columns
themes = [
    (GREEN,"🛍️","FRICTIONLESS SHOPPING ✅",
     "OOS items swap automatically.\nNo dead ends. No variant chaos.\nUser never hits a wall."),
    (AMBER,"🧠","SHOPPING BY INTENT ✅",
     "Photo→cart. YouTube link→cart.\nVoice→cart. New users get\npredictions from age+gender."),
    (BLUE,"📊","PREDICTIVE & CONFIDENT ✅",
     "Restock cart pre-built before\nuser asks. Every pick shows\nconfidence % + one-line reason."),
]
for i,(accent,icon,title,body) in enumerate(themes):
    xx = 0.35 + i*4.33
    card_box(s9,xx,0.95,4.1,1.95, accent=accent)
    txt(s9,icon,  xx+0.14,1.02,3.8,0.50, size=24, align=PP_ALIGN.CENTER)
    txt(s9,title, xx+0.14,1.54,3.8,0.28, size=10, bold=True, color=accent, align=PP_ALIGN.CENTER)
    txt(s9,body,  xx+0.14,1.84,3.8,0.98, size=9.5, color=MUTED, align=PP_ALIGN.CENTER)

# Integration strip
card_box(s9,0.35,3.02,12.6,0.40, accent=RGBColor(0x7c,0x3a,0xed))
txt(s9,"Integration-Ready:  Alexa voice  ·  Amazon Photos  ·  Engine as API  ·  Plugs into existing Amazon ecosystem",
    0.50,3.08,12.2,0.28, size=10, bold=True, color=DARK, align=PP_ALIGN.CENTER)

# Impact numbers (left)
card_box(s9,0.35,3.55,6.2,3.10, accent=GREEN)
txt(s9,"EXPECTED IMPACT", 0.50,3.60,5.9,0.22, size=8, bold=True, color=FAINT)
impacts = [
    ("⏱","26 hrs/user/year saved","10 min × 3/week. At 1M users = 26M hours returned annually."),
    ("📉","70% → <15% abandonment","Confident carts + pre-resolved OOS attack abandonment directly."),
    ("💾","99% redundant LLM calls cut","Redis hash caching: 100 queries = 1 model call + 99 cache hits."),
    ("👥","+100M newly reachable users","Voice + photo-first unlocks users excluded by text-only interfaces."),
]
for i,(icon,h,b) in enumerate(impacts):
    yy = 3.90 + i*0.68
    txt(s9,icon, 0.48,yy,0.42,0.42, size=18)
    txt(s9,h,    0.95,yy+0.02,5.35,0.26, size=11, bold=True, color=DARK)
    txt(s9,b,    0.95,yy+0.28,5.35,0.32, size=9, color=MUTED)

# Roadmap (right)
card_box(s9,6.75,3.55,5.95,3.10, accent=AMBER)
txt(s9,"EXPANSION ROADMAP", 6.90,3.60,5.6,0.22, size=8, bold=True, color=FAINT)
roadmap = [
    (GREEN,"0–3 months","Grocery · 1K users","3× faster cart · <15% abandonment"),
    (AMBER,"3–6 months","Pharmacy OTC · Multi-lang","10K users · Hindi/Tamil/Bengali"),
    (BLUE, "6–12 months","B2B + White-label API","20K+ users · licensing/cart"),
]
for i,(accent,period,title,impact) in enumerate(roadmap):
    yy = 3.95 + i*0.80
    box(s9,6.75,yy,0.08,0.72, bg_color=accent)
    txt(s9,period, 6.92,yy+0.02,1.55,0.24, size=8, bold=True, color=accent)
    txt(s9,title,  8.50,yy+0.02,4.0,0.26,  size=11, bold=True, color=DARK)
    txt(s9,impact, 8.50,yy+0.28,4.0,0.28,  size=9, color=MUTED)

# Expansion path
card_box(s9,6.75,6.12,5.95,0.44, accent=RGBColor(0x7c,0x3a,0xed))
txt(s9,"Grocery → Pharmacy → Creator → B2B → Events  |  Each step = catalog swap + prompt tweak. Zero rebuild.",
    6.90,6.17,5.65,0.32, size=9, color=DARK)

txt(s9,'Speaker note: "We hit all 3 Amazon Now areas. Engine is category-agnostic — pharmacy or creator commerce is a catalog update, not a rebuild. 1,000 grocery users in 3 months, white-label API by year-end."',
    0.35,6.95,12.6,0.45, size=8, color=FAINT, italic=True)

# ══════════════════════════════════════════════════════════════
# SLIDE 10 — CLOSING
# ══════════════════════════════════════════════════════════════
s10 = add_slide()
bg(s10, LIGHT_BG)
green_bar(s10)

# NowCart wordmark centred
txt(s10,"Now", 3.8,1.2,2.2,1.2, size=72, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)
txt(s10,"Cart  🛒", 5.7,1.2,3.8,1.2, size=72, bold=True, color=DARK, align=PP_ALIGN.CENTER)

txt(s10,'"Quick commerce solved delivery.  We solve the deciding."',
    1.5,2.50,10.3,0.50, size=16, bold=False, color=MUTED, italic=True, align=PP_ALIGN.CENTER)

divider(s10,3.18)

# 3 stats
closing_stats = [
    (GREEN_DARK,"5","Front Doors"),
    (AMBER,     "1","Reasoning Engine"),
    (ACCENT,    "0","Dead Ends"),
]
for i,(color,num,lbl) in enumerate(closing_stats):
    xx = 2.5 + i*2.8
    card_box(s10,xx,3.32,2.5,1.32)
    txt(s10,num, xx+0.1,3.38,2.3,0.72, size=48, bold=True, color=color, align=PP_ALIGN.CENTER)
    txt(s10,lbl, xx+0.1,4.08,2.3,0.42, size=12, color=MUTED, align=PP_ALIGN.CENTER)

# Deploy card
card_box(s10,2.2,4.82,8.9,1.15, accent=GREEN)
txt(s10,"Deployed on AWS  ·  Fully Functional Today",
    2.4,4.88,8.5,0.30, size=12, color=MUTED, align=PP_ALIGN.CENTER)
txt(s10,"🌐 Live App     🎥 Demo Video     💻 Source Code",
    2.4,5.18,8.5,0.35, size=13, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)
txt(s10,"Sign in as  rahul@gmail.com  (any password)  for full pantry + prediction experience",
    2.4,5.53,8.5,0.30, size=10, color=FAINT, align=PP_ALIGN.CENTER)

txt(s10,"Tell us what to make.  We assemble what to buy.",
    1.5,6.18,10.3,0.38, size=13, bold=False, color=FAINT, italic=True, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────
out = "/Users/anuj/NowCart/NowCart_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
